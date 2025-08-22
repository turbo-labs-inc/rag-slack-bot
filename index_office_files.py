#!/usr/bin/env python3
"""
Parallel document indexer for Office files from Google Drive.
Handles .docx, .xlsx, .pdf files properly.
"""

import asyncio
import time
from pathlib import Path
from typing import List, Dict, Any
from concurrent.futures import ThreadPoolExecutor
import logging
from datetime import datetime
import io

from google.oauth2 import service_account
from googleapiclient.discovery import build
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2
import chromadb
from chromadb.config import Settings
import httpx

from app.config import get_settings
from app.llm.factory import create_llm_provider

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - [Worker-%(thread)d] - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class OfficeFileIndexer:
    """Indexer for Office files from Google Drive."""
    
    def __init__(self, num_workers: int = 8):
        """Initialize with specified number of parallel workers."""
        self.num_workers = num_workers
        self.settings = get_settings()
        self.stats = {
            "total_documents": 0,
            "processed": 0,
            "failed": 0,
            "total_chunks": 0,
            "total_time": 0,
            "errors": []
        }
        
        # Initialize services
        self._init_google_drive()
        self._init_chromadb()
        self._init_llm()
        
    def _init_google_drive(self):
        """Initialize Google Drive client."""
        creds_path = Path("./credentials/google-docs-service-account.json")
        creds = service_account.Credentials.from_service_account_file(
            str(creds_path),
            scopes=['https://www.googleapis.com/auth/drive.readonly']
        )
        self.drive_service = build('drive', 'v3', credentials=creds)
        
    def _init_chromadb(self):
        """Initialize ChromaDB client."""
        self.chroma_client = chromadb.HttpClient(
            host=self.settings.chroma_host,
            port=self.settings.chroma_port
        )
        
        # Create or get collection
        try:
            self.collection = self.chroma_client.create_collection(
                name="office_documents",
                metadata={"description": "Office document embeddings"}
            )
            logger.info("Created new ChromaDB collection")
        except:
            self.collection = self.chroma_client.get_collection("office_documents")
            logger.info(f"Using existing collection with {self.collection.count()} documents")
            
    def _init_llm(self):
        """Initialize LLM provider."""
        self.llm = create_llm_provider()
        logger.info(f"Using LLM: {self.settings.llm_provider}")
        logger.info(f"Embedding model: {self.settings.ollama_embedding_model}")
        
    async def index_folder(self, folder_id: str):
        """Index all Office documents in a Google Drive folder."""
        start_time = time.time()
        
        logger.info(f"🚀 Starting parallel indexing with {self.num_workers} workers")
        logger.info(f"📁 Indexing folder: {folder_id}")
        
        # Get all documents
        documents = self._get_all_office_documents(folder_id)
        self.stats["total_documents"] = len(documents)
        
        logger.info(f"📄 Found {len(documents)} Office documents to index")
        
        # Process documents in parallel
        semaphore = asyncio.Semaphore(self.num_workers)
        
        async def process_with_semaphore(doc, index, total):
            async with semaphore:
                return await self._process_document_async(doc, index, total)
        
        # Create tasks for all documents
        tasks = [
            process_with_semaphore(doc, i, len(documents))
            for i, doc in enumerate(documents, 1)
        ]
        
        # Process all documents
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Count successes and errors
        for result in results:
            if isinstance(result, Exception):
                self.stats["failed"] += 1
                self.stats["errors"].append(str(result))
            else:
                self.stats["processed"] += 1
        
        # Calculate final stats
        self.stats["total_time"] = time.time() - start_time
        self._print_summary()
        
    def _get_all_office_documents(self, folder_id: str, recursive: bool = True) -> List[Dict]:
        """Recursively get all Office documents from folder."""
        documents = []
        
        def scan_folder(fid: str, path: str = ""):
            # Get all items in folder
            page_token = None
            while True:
                response = self.drive_service.files().list(
                    q=f"'{fid}' in parents and trashed = false",
                    fields="nextPageToken, files(id, name, mimeType)",
                    pageToken=page_token,
                    pageSize=1000
                ).execute()
                
                for file in response.get('files', []):
                    mime = file['mimeType']
                    
                    if 'folder' in mime and recursive:
                        # Recursively scan subfolder
                        scan_folder(file['id'], f"{path}/{file['name']}")
                    elif any(x in mime for x in ['wordprocessingml', 'spreadsheetml', 'presentationml', 'pdf', 'document', 'sheet', 'presentation']):
                        # Add document to list (includes PowerPoint)
                        documents.append({
                            'id': file['id'],
                            'name': file['name'],
                            'path': path,
                            'mime_type': mime
                        })
                
                page_token = response.get('nextPageToken')
                if not page_token:
                    break
        
        scan_folder(folder_id)
        return documents
        
    async def _process_document_async(self, doc: Dict, index: int, total: int):
        """Process a single document asynchronously."""
        try:
            logger.info(f"[{index}/{total}] Processing: {doc['name']}")
            
            # Download file content
            request = self.drive_service.files().get_media(fileId=doc['id'])
            file_content = io.BytesIO(request.execute())
            
            # Extract text based on file type
            text = ""
            if 'wordprocessingml' in doc['mime_type'] or doc['name'].endswith('.docx'):
                text = self._extract_docx_text(file_content)
            elif 'spreadsheetml' in doc['mime_type'] or doc['name'].endswith('.xlsx'):
                text = self._extract_xlsx_text(file_content)
            elif 'presentationml' in doc['mime_type'] or doc['name'].endswith('.pptx'):
                text = self._extract_pptx_text(file_content)
            elif 'pdf' in doc['mime_type'] or doc['name'].endswith('.pdf'):
                text = self._extract_pdf_text(file_content)
            else:
                logger.warning(f"Unsupported file type: {doc['mime_type']}")
                return
            
            if not text:
                logger.warning(f"No text extracted from {doc['name']}")
                return
            
            # Chunk the document
            chunks = self._smart_chunk(text, doc['name'])
            self.stats["total_chunks"] += len(chunks)
            
            # Generate embeddings for each chunk
            for i, chunk in enumerate(chunks):
                if chunk['text'].strip():
                    # Generate embedding
                    embedding_result = await self.llm.generate_embedding(chunk['text'])
                    
                    # Store in ChromaDB
                    self.collection.add(
                        embeddings=[embedding_result.embedding],
                        documents=[chunk['text']],
                        metadatas=[{
                            'document_id': doc['id'],
                            'document_name': doc['name'],
                            'chunk_index': i,
                            'path': doc['path']
                        }],
                        ids=[f"{doc['id']}_chunk_{i}"]
                    )
            
            logger.info(f"✅ [{index}/{total}] Indexed {len(chunks)} chunks from {doc['name']}")
            
        except Exception as e:
            logger.error(f"❌ [{index}/{total}] Error processing {doc['name']}: {e}")
            raise e
            
    def _extract_docx_text(self, file_content: io.BytesIO) -> str:
        """Extract text from Word document."""
        try:
            doc = Document(file_content)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            
            return '\n'.join(paragraphs)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
            
    def _extract_xlsx_text(self, file_content: io.BytesIO) -> str:
        """Extract text from Excel spreadsheet."""
        try:
            wb = load_workbook(file_content, read_only=True, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {e}")
            return ""
            
    def _extract_pptx_text(self, file_content: io.BytesIO) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            prs = Presentation(file_content)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {slide_num}:"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join(cell.text for cell in row.cells if cell.text)
                            if row_text.strip():
                                slide_text.append(row_text)
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_text.append(f"Notes: {notes}")
                
                if len(slide_text) > 1:  # More than just the slide number
                    text_parts.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_pdf_text(self, file_content: io.BytesIO) -> str:
        """Extract text from PDF."""
        try:
            reader = PyPDF2.PdfReader(file_content)
            text_parts = []
            
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"Page {page_num + 1}:\n{text}")
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
            
    def _smart_chunk(self, text: str, doc_name: str, chunk_size: int = 1000) -> List[Dict]:
        """Smart chunking that respects document structure."""
        chunks = []
        
        # Split by paragraphs first
        paragraphs = text.split('\n')
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_size = len(para)
            
            if current_size + para_size > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = '\n'.join(current_chunk)
                chunks.append({
                    'text': chunk_text,
                    'metadata': {'source': doc_name}
                })
                current_chunk = [para]
                current_size = para_size
            else:
                current_chunk.append(para)
                current_size += para_size
        
        # Add remaining text
        if current_chunk:
            chunks.append({
                'text': '\n'.join(current_chunk),
                'metadata': {'source': doc_name}
            })
        
        return chunks if chunks else [{'text': text, 'metadata': {'source': doc_name}}]
        
    def _print_summary(self):
        """Print indexing summary."""
        print("\n" + "="*60)
        print("📊 INDEXING COMPLETE!")
        print("="*60)
        print(f"📄 Documents processed: {self.stats['processed']}/{self.stats['total_documents']}")
        print(f"❌ Failed: {self.stats['failed']}")
        print(f"🧩 Total chunks created: {self.stats['total_chunks']}")
        print(f"⏱️  Total time: {self.stats['total_time']:.1f} seconds ({self.stats['total_time']/60:.1f} minutes)")
        
        if self.stats['processed'] > 0:
            print(f"⚡ Average time per doc: {self.stats['total_time']/self.stats['processed']:.1f}s")
            print(f"🚀 Throughput: {self.stats['total_chunks']/max(1, self.stats['total_time']):.1f} chunks/second")
        
        if self.stats['errors']:
            print(f"\n❌ Sample errors (showing first 3):")
            for error in self.stats['errors'][:3]:
                print(f"   - {error[:100]}...")
        
        print("="*60)


async def main():
    """Main entry point."""
    settings = get_settings()
    
    print("🚀 OFFICE FILE INDEXER FOR M1 ULTRA")
    print("="*60)
    
    # Configuration
    WORKERS = 8  # Start conservative, can increase to 16
    FOLDER_ID = settings.google_drive_folder_id
    
    print(f"Workers: {WORKERS}")
    print(f"Folder: {FOLDER_ID}")
    print(f"LLM: {settings.ollama_model}")
    print(f"Embeddings: {settings.ollama_embedding_model}")
    print("="*60)
    
    # Create indexer
    indexer = OfficeFileIndexer(num_workers=WORKERS)
    
    # Start indexing
    await indexer.index_folder(FOLDER_ID)
    
    print("\n💡 TIP: Increase workers to 12 or 16 for faster indexing!")
    print("   Your M1 Ultra can handle it! 💪")


if __name__ == "__main__":
    asyncio.run(main())